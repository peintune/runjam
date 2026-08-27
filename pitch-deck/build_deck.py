# -*- coding: utf-8 -*-
"""RunJam 商业 PPT 生成脚本 · 输出 16:9 PPTX"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === 品牌色 ===
C_BG          = RGBColor(0x0F, 0x12, 0x1A)
C_BG_SOFT     = RGBColor(0x1A, 0x1E, 0x2A)
C_BG_CARD     = RGBColor(0x22, 0x27, 0x36)
C_BORDER      = RGBColor(0x33, 0x39, 0x4A)
C_TEXT        = RGBColor(0xF1, 0xF5, 0xF9)
C_TEXT_DIM    = RGBColor(0x94, 0xA3, 0xB8)
C_TEXT_MUTED  = RGBColor(0x64, 0x74, 0x8B)
C_GREEN       = RGBColor(0x42, 0xB8, 0x83)
C_ORANGE      = RGBColor(0xFF, 0xA9, 0x4D)
C_RED         = RGBColor(0xCE, 0x42, 0x2B)
C_BLUE        = RGBColor(0x63, 0x66, 0xF1)
C_PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
C_AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
C_PINK        = RGBColor(0xEC, 0x48, 0x99)
C_OK          = RGBColor(0x22, 0xC5, 0x5E)
C_WARN        = RGBColor(0xF5, 0x9E, 0x0B)
C_NO          = RGBColor(0xEF, 0x44, 0x44)

def set_bg(slide, c): slide.background.fill.solid(); slide.background.fill.fore_color.rgb = c

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, corner=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if corner is not None:
        try: shp.adjustments[0] = corner
        except Exception: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None: shp.line.width = line_w
    return shp

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=C_TEXT,
             align="left", anchor="middle", font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf, m, Emu(0))
    tf.vertical_anchor = {"top":MSO_ANCHOR.TOP,"middle":MSO_ANCHOR.MIDDLE,"bottom":MSO_ANCHOR.BOTTOM}[anchor]
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left":PP_ALIGN.LEFT,"center":PP_ALIGN.CENTER,"right":PP_ALIGN.RIGHT}[align]
        r = p.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    return tb

def add_multi(slide, x, y, w, h, parts, *, align="left", anchor="top"):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf, m, Emu(0))
    tf.vertical_anchor = {"top":MSO_ANCHOR.TOP,"middle":MSO_ANCHOR.MIDDLE,"bottom":MSO_ANCHOR.BOTTOM}[anchor]
    al = {"left":PP_ALIGN.LEFT,"center":PP_ALIGN.CENTER,"right":PP_ALIGN.RIGHT}[align]
    cur = tf.paragraphs[0]; cur.alignment = al; first = True
    for part in parts:
        if part.get("newline_before") and not first: cur = tf.add_paragraph(); cur.alignment = al
        first = False
        r = cur.add_run(); r.text = part["text"]
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(part.get("size", 18))
        r.font.bold = part.get("bold", False)
        r.font.color.rgb = part.get("color", C_TEXT)
    return tb

def footer(slide, n, total):
    add_rect(slide, Inches(0.5), Inches(7.1), Inches(0.18), Inches(0.18), fill=C_GREEN, corner=0.5)
    add_text(slide, Inches(0.75), Inches(7.05), Inches(4), Inches(0.3),
             "RunJam · 一个桌面, 所有 AI Agent, 零锁定", size=10, color=C_TEXT_MUTED, anchor="middle")
    add_text(slide, Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3),
             f"{n} / {total}", size=10, color=C_TEXT_MUTED, align="right", anchor="middle")

# === 初始化 Presentation ===
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL = 12

# ============================================================
# Slide 1 — 封面
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
# 装饰色块
for x,y,w,c in [(8.5,-1.5,6.5,C_GREEN),(9.5,-1.0,5.5,C_BG),(10.3,-0.5,4.5,C_BLUE),(10.8,-0.1,4.0,C_BG)]:
    add_rect(s, Inches(x), Inches(y), Inches(w), Inches(w), fill=c, corner=0.5)
# 三色 logo
for i,(dx,dy,c) in enumerate([(1.2,1.1,C_GREEN),(1.5,1.4,C_ORANGE),(1.8,1.7,C_RED),(2.1,2.0,C_BLUE)]):
    add_rect(s, Inches(dx), Inches(dy), Inches(0.7), Inches(0.7), fill=c, corner=0.18)
add_text(s, Inches(1.2), Inches(2.7), Inches(10), Inches(0.5), "RUNJAM", size=20, bold=True, color=C_GREEN)
add_multi(s, Inches(1.2), Inches(3.2), Inches(11.5), Inches(2.4), [
    {"text":"一个桌面, ","size":56,"bold":True,"color":C_TEXT},
    {"text":"所有 AI Agent, ","size":56,"bold":True,"color":C_GREEN},
    {"text":"零锁定。","size":56,"bold":True,"color":C_ORANGE},
])
add_text(s, Inches(1.2), Inches(4.85), Inches(11.5), Inches(0.5),
         "本地优先 · 多 Agent · 多模型 · 多项目 · 统一管理", size=22, color=C_TEXT_DIM)
add_rect(s, Inches(1.2), Inches(5.7), Inches(0.08), Inches(1.2), fill=C_GREEN, corner=1.0)
add_text(s, Inches(1.45), Inches(5.7), Inches(10), Inches(0.45),
         "Claude Code  ×  Codex CLI  ×  Gemini CLI", size=18, bold=True, color=C_GREEN)
add_text(s, Inches(1.45), Inches(6.15), Inches(10.5), Inches(0.6),
         "一次配置, 所有 Agent 都能用任意模型 —— 无需 ACP 改造, 无需逐个改配置, 不上云, 不绑死。",
         size=15, color=C_TEXT_DIM)
add_text(s, Inches(1.2), Inches(6.95), Inches(11), Inches(0.3),
         "Tauri 2  ·  Rust  ·  Vue 3  ·  MIT License", size=11, color=C_TEXT_MUTED)

# ============================================================
# Slide 2 — 概览
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "01 · 概览", size=12, bold=True, color=C_GREEN)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7), "RunJam 是什么?", size=40, bold=True, color=C_TEXT)
add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.5),
         "不是 AI, 不是 IDE, 不是云端 IDE —— RunJam 是 AI Agent 桌面管理器。",
         size=18, color=C_TEXT_DIM)
add_rect(s, Inches(0.7), Inches(2.5), Inches(11.9), Inches(2.0),
         fill=C_BG_SOFT, line=C_BORDER, line_w=Pt(1), corner=0.04)
add_rect(s, Inches(0.7), Inches(2.5), Inches(0.15), Inches(2.0), fill=C_GREEN, corner=1.0)
add_multi(s, Inches(1.1), Inches(2.75), Inches(11.4), Inches(1.6), [
    {"text":"把 ","size":28,"color":C_TEXT},
    {"text":" Claude Code ","size":28,"bold":True,"color":C_ORANGE},
    {"text":"、","size":28,"color":C_TEXT},
    {"text":"Codex CLI","size":28,"bold":True,"color":C_RED},
    {"text":"、","size":28,"color":C_TEXT},
    {"text":"Gemini CLI","size":28,"bold":True,"color":C_BLUE},
    {"text":" 装进一个桌面。","size":28,"bold":True,"color":C_TEXT,"newline_before":True},
    {"text":"通过内置协议代理, ","size":22,"color":C_TEXT_DIM,"newline_before":True},
    {"text":"Anthropic ↔ OpenAI ↔ Gemini","size":22,"bold":True,"color":C_GREEN},
    {"text":" 实时互转 —— 任意 Agent 用任意模型。","size":22,"color":C_TEXT_DIM},
])
# 4 张卖点卡
cards = [("🪟","统一窗口","聊天 / 文件树 / 编辑器 / 终端, 一窗搞定",C_GREEN),
         ("🔌","任意模型","一次配置, 全 Agent 同步, 协议自动转换",C_ORANGE),
         ("🛠️","零改造","通过原生 CLI stdin/stdout 驱动 Agent",C_RED),
         ("🔒","本地优先","数据全在 ~/.runjam/, 零遥测, 零云同步",C_BLUE)]
cw,ch,gap = 2.85,2.0,0.18
x0 = (13.333-4*cw-3*gap)/2
for i,(icon,title,sub,c) in enumerate(cards):
    x = x0+i*(cw+gap)
    add_rect(s, Inches(x), Inches(5.0), Inches(cw), Inches(ch), fill=C_BG_CARD, line=C_BORDER, line_w=Pt(0.75), corner=0.08)
    add_rect(s, Inches(x), Inches(5.0), Inches(cw), Inches(0.08), fill=c)
    add_text(s, Inches(x+0.2), Inches(5.25), Inches(0.8), Inches(0.6), icon, size=28, color=c)
    add_text(s, Inches(x+0.2), Inches(5.85), Inches(cw-0.4), Inches(0.45), title, size=18, bold=True, color=C_TEXT)
    add_text(s, Inches(x+0.2), Inches(6.3), Inches(cw-0.4), Inches(0.6), sub, size=11, color=C_TEXT_DIM)
footer(s, 2, TOTAL)

# ============================================================
# Slide 3 — 8 大痛点
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "02 · 市场痛点", size=12, bold=True, color=C_RED)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "AI Agent 开发者的 8 大日常崩溃瞬间", size=36, bold=True, color=C_TEXT)
add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(0.4),
         "如果以下场景你至少中过一条, RunJam 就是为你做的。", size=16, color=C_TEXT_DIM)
pains = [("1","终端动物园","五个终端标签页, 各跑不同 Agent。Kill 错一个, 会话全丢。",C_RED),
         ("2","模型不互通","同一 prompt 测 Claude / GPT / 本地 Qwen? 三套配置挨个改。",C_ORANGE),
         ("3","协议高墙","Claude / Codex / Gemini 各讲各的协议。换模型得钻进 Agent 改。",C_AMBER),
         ("4","Token 烧钱","每个轮次重发 system prompt。哪些命中缓存? 一无所知。",C_PINK),
         ("5","环境配置地狱","新项目 / 新 Agent / 新 Key / 新 PATH / 新版本冲突。",C_PURPLE),
         ("6","单一 Agent 锁","Cursor 包了 GPT-4, Copilot 包了 OpenAI。模型换 = 工作流迁移。",C_BLUE),
         ("7","云端数据焦虑","某些 \"AI IDE\" 把你的代码 / 配置 / Prompt 发到厂商服务器。",C_RED),
         ("8","会话黑洞","关掉 IDE 会话就丢, 换台机器从头来。",C_GREEN)]
cw,ch,gx,gy = 6.1,1.05,0.18,0.18
for i,(num,title,desc,c) in enumerate(pains):
    col,row = i%2, i//2
    x,y = 0.7+col*(cw+gx), 2.2+row*(ch+gy)
    add_rect(s, Inches(x), Inches(y), Inches(cw), Inches(ch), fill=C_BG_CARD, line=C_BORDER, line_w=Pt(0.75), corner=0.05)
    add_rect(s, Inches(x+0.2), Inches(y+0.18), Inches(0.7), Inches(0.7), fill=c, corner=0.5)
    add_text(s, Inches(x+0.2), Inches(y+0.18), Inches(0.7), Inches(0.7), num, size=24, bold=True, color=C_TEXT, align="center")
    add_text(s, Inches(x+1.05), Inches(y+0.18), Inches(cw-1.25), Inches(0.4), title, size=18, bold=True, color=C_TEXT)
    add_text(s, Inches(x+1.05), Inches(y+0.6), Inches(cw-1.25), Inches(0.4), desc, size=11, color=C_TEXT_DIM)
footer(s, 3, TOTAL)

# ============================================================
# Slide 4 — 痛点 → 解法 对照表
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "03 · 解决方案", size=12, bold=True, color=C_GREEN)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "一个应用, 把 8 个痛点全部解决", size=36, bold=True, color=C_TEXT)
tx,ty,tw = 0.7, 1.7, 11.9
add_rect(s, Inches(tx), Inches(ty), Inches(tw*0.32), Inches(0.5), fill=C_BG_SOFT, line=C_BORDER, line_w=Pt(0.5))
add_rect(s, Inches(tx+tw*0.32), Inches(ty), Inches(tw*0.68), Inches(0.5), fill=C_GREEN, line=C_BORDER, line_w=Pt(0.5))
add_text(s, Inches(tx+0.2), Inches(ty), Inches(tw*0.32), Inches(0.5), "痛点", size=14, bold=True, color=C_TEXT_DIM)
add_text(s, Inches(tx+tw*0.32+0.2), Inches(ty), Inches(tw*0.68), Inches(0.5), "RunJam 的解法", size=14, bold=True, color=C_BG)
rows = [("终端动物园","统一窗口, 并行会话、侧边栏, 每个会话独立工作区",C_RED),
        ("模型不互通","统一模型中心 —— 配一次, 每个 Agent 两次点击就能绑定",C_ORANGE),
        ("协议高墙","内置协议代理, Anthropic / OpenAI / Gemini 实时互转",C_AMBER),
        ("Token 烧钱","自动检测 prompt cache + 本地响应缓存 + 每会话费用看板",C_PINK),
        ("环境配置地狱","自动检测 + 一键安装 Claude Code / Codex / Gemini CLI",C_PURPLE),
        ("单一 Agent 锁","Agent 中立 —— 换 Agent 不换工作流",C_BLUE),
        ("云端数据焦虑","本地优先, 数据全在 ~/.runjam/, API Key 进系统钥匙串, 零遥测",C_GREEN),
        ("会话黑洞","会话持久化、全文搜索、归档, 跨设备同步友好",C_AMBER)]
rh = (7.5-1.7-0.5-0.5)/len(rows)
for i,(p,sol,c) in enumerate(rows):
    y = ty+0.5+i*rh
    fr = C_BG_SOFT if i%2==0 else C_BG
    add_rect(s, Inches(tx), Inches(y), Inches(tw*0.32), Inches(rh), fill=fr, line=C_BORDER, line_w=Pt(0.3))
    add_rect(s, Inches(tx+tw*0.32), Inches(y), Inches(tw*0.68), Inches(rh), fill=fr, line=C_BORDER, line_w=Pt(0.3))
    add_rect(s, Inches(tx), Inches(y), Inches(0.08), Inches(rh), fill=c)
    add_text(s, Inches(tx+0.2), Inches(y), Inches(tw*0.32-0.3), Inches(rh), p, size=14, bold=True, color=C_TEXT)
    add_text(s, Inches(tx+tw*0.32+0.2), Inches(y), Inches(tw*0.68-0.3), Inches(rh), sol, size=12, color=C_TEXT)
footer(s, 4, TOTAL)

# ============================================================
# Slide 5 — 竞品对比矩阵 (行高已收紧, 避免越界)
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "04 · 竞争格局", size=12, bold=True, color=C_ORANGE)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "为什么不是 Cursor / Copilot / AionUI?", size=36, bold=True, color=C_TEXT)
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         "RunJam 是目前唯一同时满足「本地优先 + 任意 Agent + 协议自动转换」的产品。",
         size=14, color=C_TEXT_DIM)
headers = ["特性", "RunJam", "Cursor / Copilot", "AionUI", "原生 CLI"]
col_w = [3.6, 2.2, 2.4, 1.8, 1.9]
x0 = (13.333-sum(col_w))/2
y0 = 2.05
rh = 0.35
# 表头
xx = x0
for i,(h,w) in enumerate(zip(headers,col_w)):
    fill = C_GREEN if i==1 else C_BG_SOFT
    add_rect(s, Inches(xx), Inches(y0), Inches(w), Inches(0.5), fill=fill, line=C_BORDER, line_w=Pt(0.5))
    add_text(s, Inches(xx+0.1), Inches(y0), Inches(w-0.2), Inches(0.5), h, size=14, bold=True, color=C_BG if i==1 else C_TEXT, align="center")
    xx += w
data = [("本地优先, 不上云","✅","❌","✅","✅"),
        ("兼容任意 AI Agent CLI","✅","❌","⚠️ 仅 ACP","✅"),
        ("模型协议自动转换","✅","❌","⚠️","❌"),
        ("Agent 一键安装","✅","❌","❌","❌"),
        ("多项目并行","✅","❌","⚠️","❌"),
        ("内置编辑器 + 终端 + 文件树","✅","✅","⚠️","❌"),
        ("本地模型 (llama.cpp)","✅","❌","❌","❌"),
        ("应用管理 (自配网页应用)","✅","❌","❌","❌"),
        ("会话看板","✅","❌","❌","❌"),
        ("费用统计看板","✅","❌","❌","❌"),
        ("开源 (MIT)","✅","❌","✅","✅"),
        ("Agent 无需改造","✅","n/a","❌","n/a")]
for ri,row in enumerate(data):
    y = y0+0.5+ri*rh
    xx = x0; fr = C_BG_SOFT if ri%2==0 else C_BG
    for ci,(v,w) in enumerate(zip(row,col_w)):
        add_rect(s, Inches(xx), Inches(y), Inches(w), Inches(rh), fill=fr, line=C_BORDER, line_w=Pt(0.3))
        if ci==0:
            add_text(s, Inches(xx+0.2), Inches(y), Inches(w-0.3), Inches(rh), v, size=11, color=C_TEXT)
        else:
            col = C_OK if v=="✅" else (C_WARN if "⚠" in v else (C_NO if v=="❌" else C_TEXT_MUTED))
            add_text(s, Inches(xx), Inches(y), Inches(w), Inches(rh), v, size=13, bold=True, color=col, align="center")
        xx += w
footer(s, 5, TOTAL)

# ============================================================
# Slide 6 — 核心架构
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "05 · 架构", size=12, bold=True, color=C_BLUE)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "协议代理: RunJam 的核心杀手锏", size=36, bold=True, color=C_TEXT)
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         "Agent 通过 stdin/stdout 直接驱动 —— 协议适配在代理层完成, Agent 零改造。",
         size=14, color=C_TEXT_DIM)
cols = [("🖥️","Vue 3 前端",C_GREEN,["聊天界面","会话看板","工作区面板","设置面板"]),
        ("🤖","AI Agent",C_RED,["Claude Code","Codex CLI","Gemini CLI","(更多…)"]),
        ("🔀","RunJam 代理",C_BLUE,["协议路由","响应缓存","LLM 协议适配","Prompt cache 检测"]),
        ("☁️","LLM 接口",C_AMBER,["Anthropic","OpenAI","Google AI","DeepSeek / Qwen / 本地"])]
cw,ch,gap = 2.85,3.4,0.18
for i,(icon,title,c,items) in enumerate(cols):
    x = 0.7+i*(cw+gap)
    add_rect(s, Inches(x), Inches(2.3), Inches(cw), Inches(ch), fill=C_BG_SOFT, line=C_BORDER, line_w=Pt(0.75), corner=0.06)
    add_rect(s, Inches(x), Inches(2.3), Inches(cw), Inches(0.5), fill=c, corner=0.06)
    add_text(s, Inches(x+0.2), Inches(2.3), Inches(cw-0.4), Inches(0.5), f"{icon}  {title}", size=15, bold=True,
             color=C_BG if c in (C_GREEN,C_AMBER) else C_TEXT)
    for j,it in enumerate(items):
        add_text(s, Inches(x+0.3), Inches(3.0+j*0.45), Inches(cw-0.5), Inches(0.4), "• "+it, size=13, color=C_TEXT)
# 横向箭头
for i in range(3):
    x = 0.7+(i+1)*cw+i*gap+0.01
    add_text(s, Inches(x-0.05), Inches(4.0), Inches(gap+0.1), Inches(0.4), "→", size=22, bold=True, color=C_GREEN, align="center")
# 底部 caption
add_rect(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(1.0), fill=C_BG_CARD, line=C_GREEN, line_w=Pt(1.5), corner=0.04)
add_text(s, Inches(0.95), Inches(6.0), Inches(11.4), Inches(0.4), "Anthropic  ↔  OpenAI  ↔  Gemini", size=18, bold=True, color=C_GREEN, align="center")
add_text(s, Inches(0.95), Inches(6.4), Inches(11.4), Inches(0.5),
         "Claude Code 跑 GPT-4o · Codex 跑 Claude Sonnet · Gemini 跑 DeepSeek —— 配置一次, 全 Agent 同步。",
         size=12, color=C_TEXT_DIM, align="center")
footer(s, 6, TOTAL)

# ============================================================
# Slide 7 — 9 大功能模块
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "06 · 功能矩阵", size=12, bold=True, color=C_GREEN)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "一个桌面, 9 大功能模块", size=36, bold=True, color=C_TEXT)
feats = [("🛠️","Agent 管理","自动检测 + 一键安装 Claude / Codex / Gemini",C_GREEN),
         ("💬","统一聊天","实时流式 + 思考块 + 工具调用 + Markdown",C_ORANGE),
         ("📁","项目工作区","Monaco 编辑器 + xterm.js 终端 + 文件树",C_RED),
         ("🧠","统一模型中心","配一次, 全 Agent 同步 · 7+ 服务商预设",C_BLUE),
         ("💻","本地模型启动","llama.cpp + GGUF 一键启动, 零 API 费用",C_PURPLE),
         ("🧩","应用管理","把常用网页应用钉到 RunJam 里",C_PINK),
         ("📊","会话看板","一屏掌握所有会话状态",C_AMBER),
         ("💰","费用统计","按模型 / Agent / 天拆分 + 缓存命中率",C_GREEN),
         ("🔀","协议代理","Anthropic ↔ OpenAI ↔ Gemini 自动互转",C_BLUE)]
cw,ch,gx,gy = 4.0,1.65,0.15,0.15
for i,(icon,title,desc,c) in enumerate(feats):
    col,row = i%3, i//3
    x,y = 0.7+col*(cw+gx), 1.85+row*(ch+gy)
    add_rect(s, Inches(x), Inches(y), Inches(cw), Inches(ch), fill=C_BG_CARD, line=C_BORDER, line_w=Pt(0.75), corner=0.05)
    add_rect(s, Inches(x), Inches(y), Inches(0.1), Inches(ch), fill=c, corner=1.0)
    add_text(s, Inches(x+0.3), Inches(y+0.18), Inches(0.7), Inches(0.6), icon, size=28, color=c)
    add_text(s, Inches(x+1.05), Inches(y+0.2), Inches(cw-1.2), Inches(0.5), title, size=17, bold=True, color=C_TEXT)
    add_text(s, Inches(x+0.3), Inches(y+0.95), Inches(cw-0.5), Inches(0.6), desc, size=11, color=C_TEXT_DIM)
footer(s, 7, TOTAL)

# ============================================================
# Slide 8 — 杀手锏: 协议代理
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_rect(s, Inches(8.0), Inches(-2.0), Inches(8.0), Inches(8.0), fill=C_BG_SOFT)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "07 · 核心特性", size=12, bold=True, color=C_BLUE)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "协议代理: 为什么「一次配置」真的能成立", size=32, bold=True, color=C_TEXT)
hl = [("🔁","协议自动翻译","Anthropic ↔ OpenAI ↔ Gemini, 任意 Agent 用任意模型",C_BLUE),
      ("💾","响应缓存","重复请求本地回答, 不花 token",C_GREEN),
      ("🎯","Cache 检测","上游已缓存的 prompt, 不会再让你付一次钱",C_ORANGE),
      ("🔐","统一密钥","API Key 进系统钥匙串, 对 Agent 透明",C_PURPLE)]
for i,(icon,title,desc,c) in enumerate(hl):
    y = 1.85+i*1.05
    add_rect(s, Inches(0.7), Inches(y), Inches(0.06), Inches(0.9), fill=c, corner=1.0)
    add_text(s, Inches(0.95), Inches(y), Inches(0.7), Inches(0.9), icon, size=26, color=c)
    add_text(s, Inches(1.6), Inches(y+0.05), Inches(5.5), Inches(0.45), title, size=17, bold=True, color=C_TEXT)
    add_text(s, Inches(1.6), Inches(y+0.5), Inches(5.8), Inches(0.4), desc, size=12, color=C_TEXT_DIM)
# 右侧流程图
rx, rw = 7.5, 5.3
add_rect(s, Inches(rx), Inches(2.0), Inches(rw), Inches(0.9), fill=C_BG_CARD, line=C_ORANGE, line_w=Pt(1.5), corner=0.08)
add_text(s, Inches(rx), Inches(2.0), Inches(rw), Inches(0.9), "Claude Code  (Anthropic 协议)", size=16, bold=True, color=C_ORANGE, align="center")
add_text(s, Inches(rx+rw/2-0.3), Inches(2.95), Inches(0.6), Inches(0.3), "↓", size=20, bold=True, color=C_GREEN, align="center")
add_rect(s, Inches(rx), Inches(3.3), Inches(rw), Inches(1.0), fill=C_BLUE, corner=0.08)
add_text(s, Inches(rx), Inches(3.3), Inches(rw), Inches(1.0), "RunJam 协议代理\n实时翻译", size=15, bold=True, color=C_TEXT, align="center")
add_text(s, Inches(rx+rw/2-0.3), Inches(4.35), Inches(0.6), Inches(0.3), "↓", size=20, bold=True, color=C_GREEN, align="center")
add_rect(s, Inches(rx), Inches(4.7), Inches(rw), Inches(0.9), fill=C_BG_CARD, line=C_GREEN, line_w=Pt(1.5), corner=0.08)
add_text(s, Inches(rx), Inches(4.7), Inches(rw), Inches(0.9), "GPT-4o  (OpenAI 协议)", size=16, bold=True, color=C_GREEN, align="center")
add_text(s, Inches(rx), Inches(5.8), Inches(rw), Inches(0.4), "✅ Agent 零改造, 模型零绑定", size=13, bold=True, color=C_TEXT, align="center")
add_text(s, Inches(rx), Inches(6.2), Inches(rw), Inches(0.4), "✅ 不改 Claude Code, 也不改 GPT-4o", size=12, color=C_TEXT_DIM, align="center")
footer(s, 8, TOTAL)

# ============================================================
# Slide 9 — 三个使用场景
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "08 · 使用场景", size=12, bold=True, color=C_ORANGE)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "RunJam 的一天: 周一 / 周中 / 周五", size=36, bold=True, color=C_TEXT)
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         "三个工作日, 三种工作流, 同一个 RunJam。", size=15, color=C_TEXT_DIM)
sc = [("周一早上","三个仓库要动",
       "打开 RunJam —— 看板显示三张卡:",
       ["● runjam-core  →  Claude Code · 运行中",
        "● api-refactor  →  Codex · 空闲, 等你输入",
        "● experiments   →  Gemini · 等待审阅"],
       "点开 Codex 那张卡, 丢个新 prompt, 继续干别的。\n三个项目一个窗口, 告别切终端。", C_BLUE),
      ("周中","账单审计",
       "打开费用看板 —— 图表显示:",
       ["● 本周 60% token 烧在 api-refactor",
        "● 几乎全是 GPT-4o",
        "● 缓存命中率偏低"],
       "打开模型中心, 把 api-refactor 换成本地 Qwen-Coder 跑常规重构, 难题才留给 GPT-4o。\n下周同一张图, 费用砍半。", C_GREEN),
      ("周五","敏感任务",
       "客户发来含专有定价模型的合同条款。",
       ["● 不想让任何东西碰厂商 API",
        "● 不想上传代码, 上传 prompt",
        "● 只想用本机算力"],
       "打开本地模型启动器, 点 \"启动\" 已下好的 Qwen-72B, 全程在本机跑分析。\n数据从来没出过你的笔记本。", C_PURPLE)]
cw,ch,gx = 3.95,4.9,0.18
for i,(when,title,intro,bullets,conc,c) in enumerate(sc):
    x = 0.7+i*(cw+gx)
    add_rect(s, Inches(x), Inches(2.05), Inches(cw), Inches(ch), fill=C_BG_SOFT, line=C_BORDER, line_w=Pt(0.75), corner=0.06)
    add_rect(s, Inches(x), Inches(2.05), Inches(cw), Inches(0.7), fill=c, corner=0.06)
    add_text(s, Inches(x+0.25), Inches(2.05), Inches(cw-0.5), Inches(0.4), when, size=11, bold=True,
             color=C_BG if c in (C_GREEN,C_AMBER) else C_TEXT)
    add_text(s, Inches(x+0.25), Inches(2.35), Inches(cw-0.5), Inches(0.4), title, size=16, bold=True,
             color=C_BG if c in (C_GREEN,C_AMBER) else C_TEXT)
    add_text(s, Inches(x+0.3), Inches(3.0), Inches(cw-0.5), Inches(0.4), intro, size=12, color=C_TEXT)
    for j,b in enumerate(bullets):
        add_text(s, Inches(x+0.3), Inches(3.45+j*0.32), Inches(cw-0.5), Inches(0.3), b, size=11, color=C_TEXT)
    add_rect(s, Inches(x+0.3), Inches(4.75), Inches(cw-0.6), Inches(0.02), fill=c)
    add_text(s, Inches(x+0.3), Inches(4.9), Inches(cw-0.5), Inches(2.0), conc, size=12, color=C_TEXT, anchor="top")
footer(s, 9, TOTAL)

# ============================================================
# Slide 10 — 商业模式 + 路线图
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "09 · 商业模式 + 路线图", size=12, bold=True, color=C_GREEN)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7), "开源免费, 增值可期", size=36, bold=True, color=C_TEXT)
# 左: 商业模式
add_rect(s, Inches(0.7), Inches(1.8), Inches(5.85), Inches(5.0), fill=C_BG_SOFT, line=C_BORDER, line_w=Pt(0.75), corner=0.05)
add_rect(s, Inches(0.7), Inches(1.8), Inches(5.85), Inches(0.5), fill=C_GREEN, corner=0.05)
add_text(s, Inches(0.95), Inches(1.8), Inches(5.5), Inches(0.5), "💰  商业模式", size=16, bold=True, color=C_BG)
biz = [("开源核心 (MIT)",      "永久免费 · 个人 / 团队均可使用",  C_GREEN),
       ("云端模型按用量计费",   "走 Anthropic / OpenAI 等官方通道", C_AMBER),
       ("企业版 (未来)",       "团队协作 / 审计日志 / SSO / 私有部署", C_PURPLE),
       ("云托管 (未来)",       "可选的会话云同步 · 跨设备",          C_BLUE)]
for i,(t,d,c) in enumerate(biz):
    y = 2.5+i*0.85
    add_rect(s, Inches(0.95), Inches(y+0.1), Inches(0.3), Inches(0.3), fill=c, corner=0.5)
    add_text(s, Inches(1.4), Inches(y), Inches(4.9), Inches(0.4), t, size=15, bold=True, color=C_TEXT)
    add_text(s, Inches(1.4), Inches(y+0.4), Inches(4.9), Inches(0.4), d, size=12, color=C_TEXT_DIM)
add_text(s, Inches(0.95), Inches(6.0), Inches(5.4), Inches(0.7),
         "本地模型永远免费 —— 你只出硬件电费。\n云模型走服务商官方计费, RunJam 零抽成。",
         size=12, color=C_GREEN)
# 右: 路线图
add_rect(s, Inches(6.78), Inches(1.8), Inches(5.85), Inches(5.0), fill=C_BG_SOFT, line=C_BORDER, line_w=Pt(0.75), corner=0.05)
add_rect(s, Inches(6.78), Inches(1.8), Inches(5.85), Inches(0.5), fill=C_ORANGE, corner=0.05)
add_text(s, Inches(7.03), Inches(1.8), Inches(5.5), Inches(0.5), "🗺️  路线图", size=16, bold=True, color=C_BG)
rm = [("✅ v0.x","已完成","Agent 管理 / 统一聊天 / 多 Agent 多项目 / 模型中心 / 协议代理 / 本地模型 / 会话看板 / 应用管理",C_GREEN),
      ("🔜 v1.0","近期","Git worktree 集成 · Agent 自动更新 · 插件 / 技能系统",C_AMBER),
      ("🚧 v1.x","中期","Linux 构建 · 移动端伴侣 (只读会话视图) · 团队协作 / SSO",C_PURPLE),
      ("🔭 远期","未来","会话云同步 · 企业审计 · 私有部署 · 商业插件市场",C_BLUE)]
for i,(ver,when,what,c) in enumerate(rm):
    y = 2.5+i*0.95
    add_rect(s, Inches(6.95), Inches(y+0.05), Inches(0.12), Inches(0.85), fill=c, corner=1.0)
    add_text(s, Inches(7.2), Inches(y), Inches(1.2), Inches(0.4), ver, size=14, bold=True, color=c)
    add_text(s, Inches(8.3), Inches(y), Inches(1.0), Inches(0.4), when, size=12, color=C_TEXT_DIM)
    add_text(s, Inches(7.2), Inches(y+0.4), Inches(5.2), Inches(0.5), what, size=10, color=C_TEXT, anchor="top")
footer(s, 10, TOTAL)

# ============================================================
# Slide 11 — 加入我们
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3), "10 · 加入我们", size=12, bold=True, color=C_PINK)
add_text(s, Inches(0.7), Inches(0.85), Inches(12), Inches(0.7),
         "一个桌面, 一个社区, 一起重新定义 AI 开发", size=32, bold=True, color=C_TEXT)
add_text(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.4),
         "RunJam 100% 开源 (MIT)。我们相信工具应该属于开发者, 不属于云厂商。",
         size=15, color=C_TEXT_DIM)
cols = [("🧑‍💻","我们需要",C_GREEN,["Rust 后端工程师 (Tauri / 代理)","Vue 3 前端工程师","AI Agent 协议研究者","UI / UX 设计师","本地模型性能调优"]),
        ("⚡","现在就能上手",C_ORANGE,["⭐ Star 仓库, 让更多人看见","🐛 提 Issue / 修 Bug","📖 完善文档 / 翻译","🧪 测试 Linux 构建","🧩 贡献新 Agent 适配"]),
        ("🤝","联系方式",C_BLUE,["GitHub:  github.com/peintune/runjam","Issues:  提需求 / Bug 报告","Discussions:  功能讨论 / 反馈","PR:  欢迎提交, 见 CONTRIBUTING.md","License:  MIT"])]
cw,ch,gx = 3.95,4.0,0.18
for i,(icon,title,c,items) in enumerate(cols):
    x = 0.7+i*(cw+gx)
    add_rect(s, Inches(x), Inches(2.15), Inches(cw), Inches(ch), fill=C_BG_SOFT, line=C_BORDER, line_w=Pt(0.75), corner=0.06)
    add_rect(s, Inches(x), Inches(2.15), Inches(cw), Inches(0.55), fill=c, corner=0.06)
    add_text(s, Inches(x+0.25), Inches(2.15), Inches(cw-0.5), Inches(0.55), f"{icon}  {title}", size=15, bold=True,
             color=C_BG if c in (C_GREEN,C_AMBER) else C_TEXT)
    for j,it in enumerate(items):
        add_text(s, Inches(x+0.3), Inches(2.85+j*0.55), Inches(cw-0.5), Inches(0.5), "• "+it, size=12, color=C_TEXT)
add_text(s, Inches(0.7), Inches(6.4), Inches(11.9), Inches(0.5),
         "⭐ Star · 🐛 Issue · 🔀 PR —  让 RunJam 走得更远",
         size=18, bold=True, color=C_GREEN, align="center")
footer(s, 11, TOTAL)

# ============================================================
# Slide 12 — 收尾
# ============================================================
s = prs.slides.add_slide(BLANK); set_bg(s, C_BG)
add_rect(s, Inches(-1), Inches(0), Inches(8), Inches(7.5), fill=C_BG_SOFT)
for x,y,w,c in [(11.0,-2.0,5.0,C_GREEN),(11.5,-1.5,4.5,C_BG),(12.0,-1.0,4.0,C_ORANGE),(12.4,-0.6,3.6,C_BG)]:
    add_rect(s, Inches(x), Inches(y), Inches(w), Inches(w), fill=c, corner=0.5)
add_text(s, Inches(0.5), Inches(1.5), Inches(11), Inches(0.6), "让 AI Agent 开发, 回归你的桌面。", size=42, bold=True, color=C_TEXT, align="center")
add_text(s, Inches(0.5), Inches(2.3), Inches(11), Inches(0.5), "让选择权, 回到你手里。", size=42, bold=True, color=C_GREEN, align="center")
for dx,dy,c in [(5.4,3.5,C_GREEN),(5.8,3.8,C_ORANGE),(6.2,4.1,C_RED),(6.6,4.4,C_BLUE)]:
    add_rect(s, Inches(dx), Inches(dy), Inches(0.8), Inches(0.8), fill=c, corner=0.18)
add_text(s, Inches(0.5), Inches(5.4), Inches(11), Inches(0.4),
         "Tauri 2  ·  Rust  ·  Vue 3  ·  llama.cpp  ·  MIT License", size=14, color=C_TEXT_DIM, align="center")
add_text(s, Inches(0.5), Inches(6.0), Inches(11), Inches(0.5),
         "github.com/peintune/runjam", size=18, bold=True, color=C_GREEN, align="center")
add_text(s, Inches(0.5), Inches(6.5), Inches(11), Inches(0.4),
         "由 Rust 🦀 和 Vue 3 💚 打造", size=13, color=C_TEXT_MUTED, align="center")
add_text(s, Inches(12.6), Inches(7.1), Inches(0.6), Inches(0.3), "12 / 12", size=10, color=C_TEXT_MUTED, align="right")

# === 保存 ===
out = "/Users/guizhan/work/code/runjam/pitch-deck/RunJam-Pitch-Deck.pptx"
prs.save(out)
print(f"OK: {out}")
print(f"Slides: {len(prs.slides)}")
